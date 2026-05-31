import os
import uuid
from flask import current_app


def pdf_to_docx(input_path):
    from pdf2docx import Converter
    output_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}.docx")
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path


def pdf_to_excel(input_path):
    import pandas as pd
    import pdfplumber
    from openpyxl import Workbook  # noqa: F401 — ensures openpyxl is available for pandas

    output_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}.xlsx")
    with pdfplumber.open(input_path) as pdf:
        tables = []
        for page in pdf.pages:
            tables += page.extract_tables()
    dfs = [pd.DataFrame(table) for table in tables if table]
    df_all = pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()
    df_all.to_excel(output_path, index=False)
    return output_path


def pdf_to_csv(input_path):
    import pandas as pd
    import pdfplumber

    output_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}.csv")
    with pdfplumber.open(input_path) as pdf:
        tables = []
        for page in pdf.pages:
            tables += page.extract_tables()
    dfs = [pd.DataFrame(table) for table in tables if table]
    df_all = pd.concat(dfs, ignore_index=True) if dfs else pd.DataFrame()
    df_all.to_csv(output_path, index=False)
    return output_path


def pdf_to_image(input_path, fmt='png'):
    from pdf2image import convert_from_path

    images = convert_from_path(input_path, fmt=fmt)
    output_paths = []
    for i, img in enumerate(images, start=1):
        filename = f"{uuid.uuid4().hex}_{i}.{fmt}"
        out_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
        img.save(out_path)
        output_paths.append(out_path)
    return ','.join(output_paths)


def image_to_pdf(image_paths):
    from PIL import Image

    paths = image_paths.split(',')
    pil_images = [Image.open(p).convert('RGB') for p in paths]
    output_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}.pdf")
    pil_images[0].save(output_path, save_all=True, append_images=pil_images[1:])
    return output_path


def pdf_to_pptx(input_path):
    import tempfile
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Emu

    images = convert_from_path(input_path, dpi=150)
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    blank = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank)
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            img.save(tmp.name)
            tmp_name = tmp.name
        slide.shapes.add_picture(tmp_name, 0, 0, prs.slide_width, prs.slide_height)
        os.unlink(tmp_name)

    output_path = os.path.join(current_app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex}.pptx")
    prs.save(output_path)
    return output_path


def office_to_pdf(input_path):
    import subprocess
    upload_folder = current_app.config['UPLOAD_FOLDER']
    try:
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', upload_folder, input_path],
            capture_output=True, text=True, timeout=60
        )
    except FileNotFoundError:
        raise RuntimeError('LibreOffice is not installed. Run: sudo apt install libreoffice')
    if result.returncode != 0:
        raise RuntimeError(f'Conversion failed: {result.stderr}')
    base = os.path.splitext(os.path.basename(input_path))[0]
    out = os.path.join(upload_folder, f'{base}.pdf')
    if not os.path.exists(out):
        raise RuntimeError('Output PDF not found after conversion.')
    return out
